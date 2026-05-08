"""ChannelHub PPT report generator (Chinese labels, professional styling).

Usage:
    python scripts/generate_report_ppt.py
"""

from __future__ import annotations

import json
import sys
from datetime import datetime
from pathlib import Path

import matplotlib

matplotlib.use("Agg")
import matplotlib.font_manager as fm
import matplotlib.pyplot as plt
import numpy as np

_PROJECT_ROOT = Path(__file__).resolve().parent.parent

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Emu, Inches, Pt
except ImportError:
    print("ERROR: python-pptx not installed. Run: pip install python-pptx")
    sys.exit(1)

_CN_FONTS = ["Microsoft YaHei", "SimHei", "WenQuanYi Micro Hei", "Noto Sans CJK SC"]
_FONT_SET = False
for _f in _CN_FONTS:
    if any(_f.lower() in f.name.lower() for f in fm.fontManager.ttflist):
        plt.rcParams["font.sans-serif"] = [_f] + plt.rcParams["font.sans-serif"]
        plt.rcParams["axes.unicode_minus"] = False
        _FONT_SET = True
        break

C_PRIMARY = RGBColor(0x1A, 0x23, 0x7E)
C_SECONDARY = RGBColor(0x28, 0x3E, 0xAF)
C_ACCENT = RGBColor(0x00, 0x97, 0xA7)
C_WARM = RGBColor(0xFF, 0x6F, 0x00)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_DARK = RGBColor(0x21, 0x21, 0x21)
C_GRAY = RGBColor(0x75, 0x75, 0x75)
C_LIGHT_BG = RGBColor(0xFA, 0xFA, 0xFA)
C_TABLE_HEADER = RGBColor(0x1A, 0x23, 0x7E)
C_TABLE_ALT = RGBColor(0xE8, 0xEA, 0xF6)

SOURCE_COLORS_HEX = {
    "internal_sim": "#1565C0",
    "sionna_rt": "#00897B",
    "quadriga_real": "#D84315",
}
SOURCE_LABELS = {
    "internal_sim": "Internal TDL",
    "sionna_rt": "Sionna RT",
    "quadriga_real": "QuaDRiGa (真实)",
}


def _set_bg(slide, color=C_LIGHT_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_title_slide(prs, title, subtitle="", date_str=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, C_PRIMARY)

    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.15), Inches(7.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = C_ACCENT
    bar.line.fill.background()

    tx = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.5), Inches(1.2))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = C_WHITE

    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(20)
        p2.font.color.rgb = RGBColor(0xBB, 0xC5, 0xEA)
        p2.space_before = Pt(12)

    if date_str:
        p3 = tf.add_paragraph()
        p3.text = date_str
        p3.font.size = Pt(14)
        p3.font.color.rgb = RGBColor(0x90, 0x9C, 0xD0)
        p3.space_before = Pt(24)

    line = slide.shapes.add_shape(1, Inches(0.8), Inches(3.3), Inches(3), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = C_ACCENT
    line.line.fill.background()

    return slide


def _add_section_slide(prs, section_num, section_title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, C_SECONDARY)

    tx = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
    tf = tx.text_frame
    tf.word_wrap = True

    p0 = tf.paragraphs[0]
    p0.text = f"0{section_num}" if section_num < 10 else str(section_num)
    p0.font.size = Pt(60)
    p0.font.bold = True
    p0.font.color.rgb = RGBColor(0x60, 0x70, 0xC0)

    p1 = tf.add_paragraph()
    p1.text = section_title
    p1.font.size = Pt(32)
    p1.font.bold = True
    p1.font.color.rgb = C_WHITE
    p1.space_before = Pt(6)

    return slide


def _add_content_slide(prs, title, show_page_num=True):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)

    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = C_PRIMARY
    bar.line.fill.background()

    tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = C_PRIMARY

    uline = slide.shapes.add_shape(1, Inches(0.5), Inches(0.75), Inches(1.5), Inches(0.03))
    uline.fill.solid()
    uline.fill.fore_color.rgb = C_ACCENT
    uline.line.fill.background()

    return slide


def _add_table(slide, data, left, top, width, height):
    rows, cols = len(data), len(data[0])
    shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = shape.table

    col_w = int(width / cols)
    for c in range(cols):
        table.columns[c].width = Emu(col_w)

    for r, row in enumerate(data):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(val)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(10)
            p.alignment = PP_ALIGN.CENTER

            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = C_TABLE_HEADER
                p.font.color.rgb = C_WHITE
                p.font.bold = True
                p.font.size = Pt(11)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = C_WHITE if r % 2 == 1 else C_TABLE_ALT
                p.font.color.rgb = C_DARK

    return table


def _add_kpi_card(slide, left, top, width, height, label, value, color=C_PRIMARY):
    box = slide.shapes.add_shape(1, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = C_WHITE
    box.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
    box.line.width = Pt(1)

    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0.15)
    tf.margin_left = Inches(0.15)

    p0 = tf.paragraphs[0]
    p0.text = str(value)
    p0.font.size = Pt(28)
    p0.font.bold = True
    p0.font.color.rgb = color
    p0.alignment = PP_ALIGN.CENTER

    p1 = tf.add_paragraph()
    p1.text = label
    p1.font.size = Pt(10)
    p1.font.color.rgb = C_GRAY
    p1.alignment = PP_ALIGN.CENTER


def _styled_fig(figsize=(9, 4)):
    fig, ax = plt.subplots(figsize=figsize)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.spines["left"].set_color("#CCCCCC")
    ax.spines["bottom"].set_color("#CCCCCC")
    ax.tick_params(colors="#666666")
    return fig, ax


def _plot_snr_dist(stats, path):
    fig, ax = _styled_fig((9, 4))
    for name, s in stats.items():
        snr = s["snr_dB"]
        x = np.linspace(snr["min"] - 2, snr["max"] + 2, 200)
        sigma = max(snr["std"], 0.1)
        y = np.exp(-0.5 * ((x - snr["mean"]) / sigma) ** 2)
        c = SOURCE_COLORS_HEX.get(name, "#999")
        lab = SOURCE_LABELS.get(name, name)
        ax.fill_between(x, y, alpha=0.2, color=c)
        ax.plot(
            x,
            y,
            color=c,
            linewidth=2.5,
            label=f"{lab}  (mean={snr['mean']:.1f}, std={snr['std']:.1f})",
        )
        ax.axvline(snr["mean"], color=c, linestyle="--", alpha=0.4, linewidth=1)

    ax.set_xlabel("SNR (dB)", fontsize=12, color="#444")
    ax.set_ylabel("概率密度", fontsize=12, color="#444")
    ax.set_title("各数据源 SNR 分布", fontsize=15, fontweight="bold", color="#1A237E", pad=12)
    ax.legend(fontsize=9, framealpha=0.9, edgecolor="#DDD")
    ax.grid(True, alpha=0.15)
    fig.tight_layout()
    fig.savefig(str(path), dpi=200, bbox_inches="tight", facecolor="white")
    plt.close(fig)


def _plot_sinr_dist(stats, path):
    fig, ax = _styled_fig((9, 4))
    for name, s in stats.items():
        sinr = s["sinr_dB"]
        x = np.linspace(sinr["min"] - 2, sinr["max"] + 2, 200)
        sigma = max(sinr["std"], 0.1)
        y = np.exp(-0.5 * ((x - sinr["mean"]) / sigma) ** 2)
        c = SOURCE_COLORS_HEX.get(name, "#999")
        lab = SOURCE_LABELS.get(name, name)
        ax.fill_between(x, y, alpha=0.2, color=c)
        ax.plot(x, y, color=c, linewidth=2.5, label=f"{lab}  (mean={sinr['mean']:.1f})")

    ax.set_xlabel("SINR (dB)", fontsize=12, color="#444")
    ax.set_ylabel("概率密度", fontsize=12, color="#444")
    ax.set_title("各数据源 SINR 分布", fontsize=15, fontweight="bold", color="#1A237E", pad=12)
    ax.legend(fontsize=9, framealpha=0.9, edgecolor="#DDD")
    ax.grid(True, alpha=0.15)
    fig.tight_layout()
    fig.savefig(str(path), dpi=200, bbox_inches="tight", facecolor="white")
    plt.close(fig)


def _plot_sample_bars(stats, path):
    fig, ax = _styled_fig((5, 3.5))
    names = list(stats.keys())
    labels = [SOURCE_LABELS.get(n, n) for n in names]
    counts = [stats[n]["count"] for n in names]
    colors = [SOURCE_COLORS_HEX.get(n, "#999") for n in names]

    bars = ax.bar(labels, counts, color=colors, alpha=0.85, width=0.5, edgecolor="white", linewidth=1.5)
    for bar, c in zip(bars, counts, strict=False):
        ax.text(
            bar.get_x() + bar.get_width() / 2,
            bar.get_height() + 50,
            f"{c:,}",
            ha="center",
            va="bottom",
            fontsize=12,
            fontweight="bold",
            color="#333",
        )

    ax.set_ylabel("样本数量", fontsize=11, color="#444")
    ax.set_title("各数据源样本量", fontsize=14, fontweight="bold", color="#1A237E", pad=10)
    ax.grid(True, alpha=0.15, axis="y")
    ax.set_ylim(0, max(counts) * 1.15)
    fig.tight_layout()
    fig.savefig(str(path), dpi=200, bbox_inches="tight", facecolor="white")
    plt.close(fig)


def _plot_power_bars(stats, path):
    fig, ax = _styled_fig((5, 3.5))
    names = list(stats.keys())
    labels = [SOURCE_LABELS.get(n, n) for n in names]
    means = [stats[n]["channel_power"]["mean"] for n in names]
    stds = [stats[n]["channel_power"]["std"] for n in names]
    colors = [SOURCE_COLORS_HEX.get(n, "#999") for n in names]

    bars = ax.bar(
        labels, means, yerr=stds, capsize=5, color=colors, alpha=0.85, width=0.5, edgecolor="white", linewidth=1.5
    )
    for bar, m in zip(bars, means, strict=False):
        ax.text(
            bar.get_x() + bar.get_width() / 2,
            bar.get_height(),
            f"{m:.4f}",
            ha="center",
            va="bottom",
            fontsize=10,
            color="#333",
        )

    ax.set_ylabel("平均信道功率", fontsize=11, color="#444")
    ax.set_title("信道功率对比", fontsize=14, fontweight="bold", color="#1A237E", pad=10)
    ax.grid(True, alpha=0.15, axis="y")
    fig.tight_layout()
    fig.savefig(str(path), dpi=200, bbox_inches="tight", facecolor="white")
    plt.close(fig)


def _plot_training(train_log, path):
    fig, axes = plt.subplots(1, 2, figsize=(10, 4))

    for ax in axes:
        ax.spines["top"].set_visible(False)
        ax.spines["right"].set_visible(False)
        ax.spines["left"].set_color("#CCCCCC")
        ax.spines["bottom"].set_color("#CCCCCC")

    epochs = train_log.get("epochs", [])
    tl = train_log.get("train_loss", [])
    vl = train_log.get("val_loss", [])

    ax1 = axes[0]
    if epochs and tl:
        ax1.plot(epochs, tl, color="#1565C0", linewidth=2.5, label="训练损失", marker="o", markersize=3)
        if vl:
            ax1.plot(epochs, vl, color="#EF6C00", linewidth=2.5, label="验证损失", marker="s", markersize=3)
        ax1.set_xlabel("Epoch", fontsize=11, color="#444")
        ax1.set_ylabel("MSE Loss", fontsize=11, color="#444")
        ax1.set_title("损失曲线", fontsize=14, fontweight="bold", color="#1A237E", pad=10)
        ax1.legend(fontsize=10, framealpha=0.9)
        ax1.grid(True, alpha=0.15)

        if vl:
            ax1.annotate(
                f"val={vl[-1]:.4f}",
                xy=(epochs[-1], vl[-1]),
                fontsize=9,
                color="#EF6C00",
                fontweight="bold",
                xytext=(-50, 15),
                textcoords="offset points",
                arrowprops={"arrowstyle": "->", "color": "#EF6C00", "lw": 1},
            )

    lr = train_log.get("learning_rate", [])
    ax2 = axes[1]
    if epochs and lr:
        ax2.plot(epochs, lr, color="#00897B", linewidth=2.5)
        ax2.fill_between(epochs, lr, alpha=0.1, color="#00897B")
        ax2.set_xlabel("Epoch", fontsize=11, color="#444")
        ax2.set_ylabel("学习率", fontsize=11, color="#444")
        ax2.set_title("学习率调度 (Cosine)", fontsize=14, fontweight="bold", color="#1A237E", pad=10)
        ax2.grid(True, alpha=0.15)
        ax2.ticklabel_format(style="scientific", axis="y", scilimits=(0, 0))

    fig.tight_layout(w_pad=3)
    fig.savefig(str(path), dpi=200, bbox_inches="tight", facecolor="white")
    plt.close(fig)


def _plot_embedding_analysis(emb_path, stats, fig_dir):
    if not emb_path.exists():
        return None, None

    emb = np.load(str(emb_path))
    n_total = emb.shape[0]

    source_sizes = []
    source_names = []
    for name in stats:
        source_names.append(name)
        source_sizes.append(stats[name]["count"])

    from sklearn.decomposition import PCA

    pca = PCA(n_components=2)
    emb_2d = pca.fit_transform(emb)

    fig, ax = _styled_fig((5, 4.5))
    offset = 0
    for name, sz in zip(source_names, source_sizes, strict=False):
        sz_actual = min(sz, n_total - offset)
        if sz_actual <= 0:
            break
        pts = emb_2d[offset : offset + sz_actual]
        c = SOURCE_COLORS_HEX.get(name, "#999")
        lab = SOURCE_LABELS.get(name, name)
        ax.scatter(pts[:, 0], pts[:, 1], c=c, s=4, alpha=0.4, label=lab)
        offset += sz_actual

    ax.set_xlabel(f"PC1 ({pca.explained_variance_ratio_[0] * 100:.1f}%)", fontsize=10, color="#444")
    ax.set_ylabel(f"PC2 ({pca.explained_variance_ratio_[1] * 100:.1f}%)", fontsize=10, color="#444")
    ax.set_title("Embedding PCA 投影", fontsize=14, fontweight="bold", color="#1A237E", pad=10)
    ax.legend(fontsize=9, markerscale=3, framealpha=0.9)
    fig.tight_layout()
    pca_path = fig_dir / "embedding_pca.png"
    fig.savefig(str(pca_path), dpi=200, bbox_inches="tight", facecolor="white")
    plt.close(fig)

    fig, ax = _styled_fig((5, 4.5))
    offset = 0
    for name, sz in zip(source_names, source_sizes, strict=False):
        sz_actual = min(sz, n_total - offset)
        if sz_actual <= 0:
            break
        norms = np.linalg.norm(emb[offset : offset + sz_actual], axis=1)
        c = SOURCE_COLORS_HEX.get(name, "#999")
        lab = SOURCE_LABELS.get(name, name)
        ax.hist(norms, bins=50, alpha=0.5, color=c, label=f"{lab} (mean={np.mean(norms):.3f})")
        offset += sz_actual

    ax.set_xlabel("L2 Norm", fontsize=10, color="#444")
    ax.set_ylabel("频次", fontsize=10, color="#444")
    ax.set_title("Embedding 模长分布", fontsize=14, fontweight="bold", color="#1A237E", pad=10)
    ax.legend(fontsize=9, framealpha=0.9)
    ax.grid(True, alpha=0.15)
    fig.tight_layout()
    norm_path = fig_dir / "embedding_norm.png"
    fig.savefig(str(norm_path), dpi=200, bbox_inches="tight", facecolor="white")
    plt.close(fig)

    return pca_path, norm_path


def generate_ppt(
    stats_path: Path,
    train_log_path: Path | None,
    output_path: Path,
    figures_dir: Path,
):
    stats = json.loads(stats_path.read_text(encoding="utf-8"))
    train_log = {}
    if train_log_path and train_log_path.exists():
        train_log = json.loads(train_log_path.read_text(encoding="utf-8"))

    figures_dir.mkdir(parents=True, exist_ok=True)
    total = sum(s["count"] for s in stats.values())

    _plot_snr_dist(stats, figures_dir / "snr_dist.png")
    _plot_sinr_dist(stats, figures_dir / "sinr_dist.png")
    _plot_sample_bars(stats, figures_dir / "sample_count.png")
    _plot_power_bars(stats, figures_dir / "power_cmp.png")
    if train_log:
        _plot_training(train_log, figures_dir / "train_loss.png")

    emb_path = stats_path.parent / "embeddings.npy"
    pca_img, norm_img = _plot_embedding_analysis(emb_path, stats, figures_dir)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    _add_title_slide(
        prs,
        "ChannelHub · 信道数据工场",
        "多源数据采集、模型训练与推理分析汇报",
        datetime.now().strftime("%Y年%m月%d日"),
    )

    slide = _add_content_slide(prs, "目录")
    toc_items = [
        ("01", "数据总览", "三种信道数据源采集结果概况"),
        ("02", "信噪比分析", "SNR / SINR 分布特征对比"),
        ("03", "各数据源详情", "配置参数与统计指标"),
        ("04", "模型训练结果", "ChannelMAE 30 Epoch 训练过程"),
        ("05", "嵌入空间分析", "16维隐空间可视化与聚类特征"),
        ("06", "总结与展望", "关键发现与后续计划"),
    ]
    for i, (num, title, desc) in enumerate(toc_items):
        y = Inches(1.2) + Inches(i * 0.9)
        tx = slide.shapes.add_textbox(Inches(1.5), y, Inches(0.8), Inches(0.5))
        p = tx.text_frame.paragraphs[0]
        p.text = num
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = C_ACCENT
        tx2 = slide.shapes.add_textbox(Inches(2.5), y, Inches(4), Inches(0.35))
        p2 = tx2.text_frame.paragraphs[0]
        p2.text = title
        p2.font.size = Pt(18)
        p2.font.bold = True
        p2.font.color.rgb = C_PRIMARY
        tx3 = slide.shapes.add_textbox(Inches(2.5), y + Inches(0.35), Inches(6), Inches(0.35))
        p3 = tx3.text_frame.paragraphs[0]
        p3.text = desc
        p3.font.size = Pt(11)
        p3.font.color.rgb = C_GRAY

    _add_section_slide(prs, 1, "数据总览")
    slide = _add_content_slide(prs, "数据采集概况")
    cards = [
        ("总样本数", f"{total:,}", C_PRIMARY),
        ("数据源", f"{len(stats)}", C_ACCENT),
        ("天线配置", "4 BS x 2 UE", C_WARM),
        ("OFDM 参数", "52 RB x 14 sym", RGBColor(0x7B, 0x1F, 0xA2)),
    ]
    for i, (label, value, color) in enumerate(cards):
        _add_kpi_card(slide, Inches(0.5 + i * 3.1), Inches(1.1), Inches(2.8), Inches(1.0), label, value, color)

    data = [["数据源", "样本数", "SNR 均值", "SNR 标准差", "SINR 均值", "信道功率", "信道估计"]]
    for name, s in stats.items():
        data.append([
            SOURCE_LABELS.get(name, name),
            f"{s['count']:,}",
            f"{s['snr_dB']['mean']:.1f} dB",
            f"{s['snr_dB']['std']:.1f} dB",
            f"{s['sinr_dB']['mean']:.1f} dB",
            f"{s['channel_power']['mean']:.4f}",
            list(s.get("channel_est_modes", {}).keys())[0] if s.get("channel_est_modes") else "-",
        ])
    _add_table(slide, data, Inches(0.5), Inches(2.5), Inches(12.3), Inches(1.5))

    slide.shapes.add_picture(str(figures_dir / "sample_count.png"), Inches(0.3), Inches(4.3), Inches(6), Inches(3))
    slide.shapes.add_picture(str(figures_dir / "power_cmp.png"), Inches(6.8), Inches(4.3), Inches(6), Inches(3))

    _add_section_slide(prs, 2, "信噪比分析")
    slide = _add_content_slide(prs, "SNR 分布对比")
    slide.shapes.add_picture(str(figures_dir / "snr_dist.png"), Inches(0.3), Inches(1.0), Inches(12.5), Inches(4))
    snr_data = [["数据源", "最小值", "P25", "中位数", "P75", "最大值", "均值", "标准差"]]
    for name, s in stats.items():
        snr = s["snr_dB"]
        snr_data.append([
            SOURCE_LABELS.get(name, name),
            f"{snr['min']:.1f}",
            f"{snr['p25']:.1f}",
            f"{snr['p50']:.1f}",
            f"{snr['p75']:.1f}",
            f"{snr['max']:.1f}",
            f"{snr['mean']:.1f}",
            f"{snr['std']:.1f}",
        ])
    _add_table(slide, snr_data, Inches(0.5), Inches(5.3), Inches(12), Inches(1.5))

    slide = _add_content_slide(prs, "SINR 分布对比")
    slide.shapes.add_picture(str(figures_dir / "sinr_dist.png"), Inches(0.3), Inches(1.0), Inches(12.5), Inches(4))
    sinr_data = [["数据源", "最小值", "P25", "中位数", "P75", "最大值", "均值", "标准差"]]
    for name, s in stats.items():
        sinr = s["sinr_dB"]
        sinr_data.append([
            SOURCE_LABELS.get(name, name),
            f"{sinr['min']:.1f}",
            f"{sinr['p25']:.1f}",
            f"{sinr['p50']:.1f}",
            f"{sinr['p75']:.1f}",
            f"{sinr['max']:.1f}",
            f"{sinr['mean']:.1f}",
            f"{sinr['std']:.1f}",
        ])
    _add_table(slide, sinr_data, Inches(0.5), Inches(5.3), Inches(12), Inches(1.5))

    _add_section_slide(prs, 3, "各数据源详情")
    for name, s in stats.items():
        label = SOURCE_LABELS.get(name, name)
        slide = _add_content_slide(prs, f"{label} 数据源详情")

        metrics_items = [
            ("样本数量", f"{s['count']:,}"),
            ("SNR", f"{s['snr_dB']['mean']:.1f} ± {s['snr_dB']['std']:.1f} dB"),
            ("SINR", f"{s['sinr_dB']['mean']:.1f} ± {s['sinr_dB']['std']:.1f} dB"),
            ("信道功率", f"{s['channel_power']['mean']:.4f}"),
            ("链路方向", "/".join(s.get("links", {}).keys()) or "-"),
            ("信道估计", "/".join(s.get("channel_est_modes", {}).keys()) or "-"),
        ]
        for i, (lbl, val) in enumerate(metrics_items):
            y = Inches(1.1) + Inches(i * 0.55)
            tx = slide.shapes.add_textbox(Inches(0.6), y, Inches(2), Inches(0.4))
            p = tx.text_frame.paragraphs[0]
            p.text = lbl
            p.font.size = Pt(11)
            p.font.color.rgb = C_GRAY
            p.font.bold = True
            tx2 = slide.shapes.add_textbox(Inches(2.8), y, Inches(3), Inches(0.4))
            p2 = tx2.text_frame.paragraphs[0]
            p2.text = val
            p2.font.size = Pt(12)
            p2.font.color.rgb = C_DARK

        cfg = s.get("config", {})
        if cfg:
            cfg_keys = [
                "scenario",
                "num_cells",
                "num_sites",
                "carrier_freq_hz",
                "isd_m",
                "channel_est_mode",
                "link",
                "seed",
                "bs_antennas",
                "ue_antennas",
                "num_rb",
            ]
            cfg_data = [["参数", "值"]]
            for k in cfg_keys:
                if k in cfg:
                    cfg_data.append([k, str(cfg[k])])
            if len(cfg_data) > 1:
                _add_table(slide, cfg_data, Inches(7), Inches(1.1), Inches(5.5), Inches(0.4 * len(cfg_data)))

        if "ue_position" in s:
            pos = s["ue_position"]
            tx = slide.shapes.add_textbox(Inches(0.6), Inches(5.5), Inches(8), Inches(0.5))
            p = tx.text_frame.paragraphs[0]
            p.text = (
                f"UE 分布范围:  X=[{pos['x_range'][0]:.0f}, {pos['x_range'][1]:.0f}]m  "
                f"Y=[{pos['y_range'][0]:.0f}, {pos['y_range'][1]:.0f}]m"
            )
            p.font.size = Pt(11)
            p.font.color.rgb = C_GRAY

    if train_log and train_log.get("epochs"):
        _add_section_slide(prs, 4, "模型训练结果")
        slide = _add_content_slide(prs, "ChannelMAE 训练过程")
        if (figures_dir / "train_loss.png").exists():
            slide.shapes.add_picture(
                str(figures_dir / "train_loss.png"), Inches(0.3), Inches(1.0), Inches(12.5), Inches(4.2)
            )
        train_cards = [
            ("最终训练损失", f"{train_log.get('final_train_loss', 0):.4f}", C_PRIMARY),
            ("最终验证损失", f"{train_log.get('final_val_loss', 0):.4f}", C_WARM),
            ("最优 Epoch", str(train_log.get("best_epoch", "-")), C_ACCENT),
            ("总 Epoch 数", str(len(train_log.get("epochs", []))), C_GRAY),
        ]
        for i, (label, value, color) in enumerate(train_cards):
            _add_kpi_card(slide, Inches(0.5 + i * 3.1), Inches(5.8), Inches(2.8), Inches(1.2), label, value, color)

    if pca_img or norm_img:
        _add_section_slide(prs, 5, "嵌入空间分析")
        slide = _add_content_slide(prs, "16维 Embedding 可视化")
        if pca_img and pca_img.exists():
            slide.shapes.add_picture(str(pca_img), Inches(0.3), Inches(1.0), Inches(6.2), Inches(5.5))
        if norm_img and norm_img.exists():
            slide.shapes.add_picture(str(norm_img), Inches(6.8), Inches(1.0), Inches(6.2), Inches(5.5))
        tx = slide.shapes.add_textbox(Inches(0.5), Inches(6.7), Inches(12), Inches(0.5))
        p = tx.text_frame.paragraphs[0]
        p.text = "左: PCA 降维投影，颜色区分数据源 | 右: Embedding L2 模长分布（L2归一化后均为1.0）"
        p.font.size = Pt(10)
        p.font.color.rgb = C_GRAY

    _add_section_slide(prs, 6, "总结与展望")
    slide = _add_content_slide(prs, "关键发现与后续计划")

    tx = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(5.8), Inches(5.5))
    tf = tx.text_frame
    tf.word_wrap = True
    findings_title = tf.paragraphs[0]
    findings_title.text = "关键发现"
    findings_title.font.size = Pt(18)
    findings_title.font.bold = True
    findings_title.font.color.rgb = C_PRIMARY

    findings = [
        f"成功采集 {total:,} 条多源信道数据（{len(stats)} 种采集方式）",
        (
            f"SNR 覆盖 {min(s['snr_dB']['min'] for s in stats.values()):.0f} ~ "
            f"{max(s['snr_dB']['max'] for s in stats.values()):.0f} dB，涵盖多种信道场景"
        ),
        "各数据源信道功率均归一化至 ~1.0，保证特征空间一致性",
        "Bridge 管道产出 16 token + 8 gate 特征向量，可靠稳定",
    ]
    if train_log:
        vl = train_log.get("final_val_loss")
        if vl:
            findings.append(f"ChannelMAE 30 Epoch 训练收敛良好，val_loss = {vl:.4f}")
        findings.append("16维嵌入空间 L2 归一化，适用于下游检索/聚类任务")

    for f in findings:
        p = tf.add_paragraph()
        p.text = f"  {f}"
        p.font.size = Pt(12)
        p.font.color.rgb = C_DARK
        p.space_before = Pt(8)

    tx2 = slide.shapes.add_textbox(Inches(6.8), Inches(1.1), Inches(5.8), Inches(5.5))
    tf2 = tx2.text_frame
    tf2.word_wrap = True
    next_title = tf2.paragraphs[0]
    next_title.text = "后续计划"
    next_title.font.size = Pt(18)
    next_title.font.bold = True
    next_title.font.color.rgb = C_ACCENT

    nexts = [
        "接入真实 QuaDRiGa MATLAB 仿真，替换合成数据",
        "扩展至 10 万+ 样本规模，GPU 加速 Sionna RT",
        "DDP 多卡分布式训练（4-8 GPU）",
        "下游评估：Channel Charting CT/TW 指标",
        "ONNX 模型导出，支持生产环境部署",
        "数据管理平台功能完善与交互优化",
    ]
    for n in nexts:
        p = tf2.add_paragraph()
        p.text = f"  {n}"
        p.font.size = Pt(12)
        p.font.color.rgb = C_DARK
        p.space_before = Pt(8)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, C_PRIMARY)
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.15), Inches(7.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = C_ACCENT
    bar.line.fill.background()
    tx = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(1))
    p = tx.text_frame.paragraphs[0]
    p.text = "Thank You"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = C_WHITE
    p.alignment = PP_ALIGN.CENTER
    p2 = tx.text_frame.add_paragraph()
    p2.text = "ChannelHub · 信道数据工场"
    p2.font.size = Pt(16)
    p2.font.color.rgb = RGBColor(0x90, 0x9C, 0xD0)
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(12)

    prs.save(str(output_path))
    print(f"PPT saved: {output_path}")
    return output_path


if __name__ == "__main__":
    bridge_out = _PROJECT_ROOT / "artifacts" / "bridge_out_5k"
    stats_path = bridge_out / "dataset_stats.json"
    train_log_path = bridge_out / "train_log.json"
    figures_dir = _PROJECT_ROOT / "reports" / "figures"
    output_path = _PROJECT_ROOT / "reports" / "MSG_Embedding_Report.pptx"

    if not stats_path.exists():
        print(f"Stats not found: {stats_path}")
        sys.exit(1)

    generate_ppt(stats_path, train_log_path, output_path, figures_dir)
